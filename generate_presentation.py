from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont


ASSETS_DIR = Path("assets")


def ensure_assets():
    ASSETS_DIR.mkdir(exist_ok=True)
    create_cafe_background(ASSETS_DIR / "cafe_background.jpg")
    create_college_logo(ASSETS_DIR / "college_logo.png")
    create_industry_graphic(ASSETS_DIR / "industry_graphic.png")
    create_comparison_chart(ASSETS_DIR / "comparison_chart.png")
    create_busy_cafe_visual(ASSETS_DIR / "busy_cafe.png")
    create_objective_icon(ASSETS_DIR / "objectives.png")


def base_canvas(width=1280, height=720, color="#f5e6d3"):
    return Image.new("RGB", (width, height), color)


def load_font(size: int):
    try:
        return ImageFont.truetype("DejaVuSans-Bold.ttf", size)
    except OSError:
        return ImageFont.load_default()


def create_cafe_background(path: Path):
    if path.exists():
        return
    img = base_canvas(color="#d7b790")
    draw = ImageDraw.Draw(img)
    cup_color = "#704214"
    saucer_color = "#f5eadf"
    draw.rectangle([0, img.height - 220, img.width, img.height], fill="#a36d3c")
    cup_bbox = [img.width // 2 - 140, img.height - 380, img.width // 2 + 140, img.height - 180]
    draw.ellipse(cup_bbox, fill=cup_color)
    draw.rectangle([img.width // 2 - 90, img.height - 460, img.width // 2 + 90, img.height - 380], fill=cup_color)
    draw.pieslice([img.width // 2 - 100, img.height - 500, img.width // 2 + 100, img.height - 300], 0, 180, fill=saucer_color)
    steam_color = "#ffffff"
    for offset in (-60, 0, 60):
        draw.line(
            [
                (img.width // 2 + offset, img.height - 480),
                (img.width // 2 + offset - 20, img.height - 620),
            ],
            fill=steam_color,
            width=6,
        )
    font = load_font(90)
    draw.text((60, 60), "Smart Café", fill="#3b2512", font=font)
    img.save(path, quality=95)


def create_college_logo(path: Path):
    if path.exists():
        return
    img = Image.new("RGBA", (512, 512), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    draw.ellipse([16, 16, 496, 496], fill="#004d99")
    draw.ellipse([96, 96, 416, 416], fill="#ffffff")
    draw.polygon([(256, 120), (360, 392), (152, 392)], fill="#ffb400")
    draw.text((205, 220), "SC", fill="#004d99", font=load_font(64))
    img.save(path)


def create_industry_graphic(path: Path):
    if path.exists():
        return
    img = base_canvas(color="#fdf8f3")
    draw = ImageDraw.Draw(img)
    bars = [220, 340, 480, 620, 700]
    for idx, height in enumerate(bars):
        x0 = 140 + idx * 160
        draw.rectangle([x0, img.height - height, x0 + 100, img.height - 120], fill="#b26b2a")
    draw.line([120, img.height - 120, 1100, img.height - 120], fill="#573116", width=6)
    draw.ellipse([950, 80, 1080, 210], fill="#f2c94c")
    draw.text((980, 120), "₹9B", fill="#573116", font=load_font(48))
    draw.text((120, 60), "Café Industry Growth", fill="#573116", font=load_font(64))
    img.save(path, quality=95)


def create_comparison_chart(path: Path):
    if path.exists():
        return
    img = base_canvas(color="#fff9f1")
    draw = ImageDraw.Draw(img)
    columns = ["Features", "Cost", "Localization"]
    for idx, title in enumerate(columns):
        draw.rectangle([60 + idx * 380, 80, 400 + idx * 380, 140], fill="#704214")
        draw.text((80 + idx * 380, 92), title, fill="#ffffff", font=load_font(32))
    rows = ["Legacy POS", "Premium Suites", "Smart Café"]
    for ridx, row in enumerate(rows):
        draw.rectangle([60, 160 + ridx * 160, 1180, 300 + ridx * 160], outline="#704214", width=4)
        draw.text((80, 200 + ridx * 160), row, fill="#3b2512", font=load_font(36))
    draw.text((500, 220), "High\nCapex", fill="#bf360c", font=load_font(32))
    draw.text((880, 360), "Limited\nSupport", fill="#bf360c", font=load_font(32))
    draw.text((520, 520), "Affordable,\nPay-as-you-go", fill="#2e7d32", font=load_font(32))
    draw.text((900, 520), "Regional\nLanguages", fill="#2e7d32", font=load_font(32))
    img.save(path, quality=95)


def create_busy_cafe_visual(path: Path):
    if path.exists():
        return
    img = base_canvas(color="#fbeee2")
    draw = ImageDraw.Draw(img)
    draw.rectangle([80, 200, 1200, 520], fill="#d7b58c")
    for x in range(0, 6):
        table_x = 120 + x * 180
        draw.ellipse([table_x, 320, table_x + 120, 440], fill="#8d6e63")
    draw.rounded_rectangle([500, 80, 780, 200], radius=40, fill="#4e342e")
    draw.text((530, 120), "Manual Billing", fill="#ffffff", font=load_font(40))
    draw.rounded_rectangle([900, 80, 1180, 200], radius=40, fill="#bf360c")
    draw.text((930, 120), "Order Queue", fill="#ffffff", font=load_font(40))
    draw.rounded_rectangle([80, 80, 360, 200], radius=40, fill="#3e2723")
    draw.text((120, 120), "Kitchen", fill="#ffffff", font=load_font(40))
    draw.text((120, 580), "Inefficient coordination & missing records", fill="#3b2512", font=load_font(36))
    img.save(path, quality=95)


def create_objective_icon(path: Path):
    if path.exists():
        return
    img = base_canvas(color="#fef6eb")
    draw = ImageDraw.Draw(img)
    draw.ellipse([120, 120, 520, 520], fill="#ffe0b2", outline="#ff9800", width=10)
    draw.rectangle([320, 180, 600, 500], fill="#ffcc80")
    draw.rectangle([280, 220, 560, 360], fill="#fff3e0")
    draw.rectangle([280, 380, 560, 440], fill="#fff3e0")
    draw.line([320, 420, 520, 420], fill="#ff9800", width=12)
    draw.line([320, 460, 520, 460], fill="#ff9800", width=12)
    draw.line([360, 250, 460, 330], fill="#ff9800", width=12)
    draw.line([460, 330, 520, 260], fill="#2e7d32", width=12)
    draw.text((200, 560), "Automation Goals", fill="#bf360c", font=load_font(40))
    img.save(path, quality=95)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(
        str(ASSETS_DIR / "cafe_background.jpg"),
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(2.2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.text = "Smart Cafe Management System"
    title_para = title_tf.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(6.5), Inches(1.2))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = "B.Tech Mini Project"
    subtitle_para = subtitle_tf.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

    presenter_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(6.5), Inches(1.5))
    presenter_tf = presenter_box.text_frame
    presenter_tf.text = "Presenter: Aryan Sharma (Roll No. BT21CS045)\nDepartment of Computer Science & Engineering"
    for paragraph in presenter_tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    logo_path = ASSETS_DIR / "college_logo.png"
    slide.shapes.add_picture(str(logo_path), Inches(9.5), Inches(0.6), height=Inches(1.6))


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], image_path: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(40)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(6.5), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            tf.text = bullet
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
            para.text = bullet
        para.level = 0
        para.font.size = Pt(24)

    if image_path:
        slide.shapes.add_picture(str(image_path), Inches(7.5), Inches(2), width=Inches(3.8))

    return slide


def add_section_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(42)
    slide.placeholders[1].text = subtitle
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(28)
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_presentation():
    ensure_assets()
    prs = Presentation()

    add_title_slide(prs)

    slides_data = [
        (
            "Motivation & Background",
            [
                "Digitize daily café operations to reduce manual overhead and billing errors.",
                "Empower small Indian cafés with technology tailored to their budgets and workflows.",
                "Ride the wave of digital payments, UPI adoption, and on-demand ordering trends.",
            ],
            ASSETS_DIR / "industry_graphic.png",
        ),
        (
            "Literature Review / Existing Systems",
            [
                "Evaluated POS offerings such as Toast, Square, and legacy desktop billing suites.",
                "Identified pain points: high licensing fees, hardware lock-ins, complicated training.",
                "Opportunity for a light-weight, browser-based system with local language support.",
            ],
            ASSETS_DIR / "comparison_chart.png",
        ),
        (
            "Problem Statement",
            [
                "Manual billing produces pricing errors and inconsistent receipts.",
                "Kitchen receives delayed or unclear order tickets, causing slow service.",
                "No central view of inventory consumption or wastage trends.",
                "Customer preferences and loyalty data rarely captured or analyzed.",
            ],
            ASSETS_DIR / "busy_cafe.png",
        ),
        (
            "Objectives",
            [
                "Automate order capture, billing, and payment reconciliation in real time.",
                "Provide live inventory tracking with automated low-stock alerts.",
                "Enable cashier–kitchen coordination via synchronized order screens.",
                "Build customer engagement with digital receipts and loyalty tracking.",
            ],
            ASSETS_DIR / "objectives.png",
        ),
        (
            "Project Scope",
            [
                "Develop a responsive web POS optimized for tablets and kiosks.",
                "Implement admin dashboard for menu, pricing, and staff management.",
                "Integrate analytics for daily sales, category performance, and peak hours.",
                "Offer configurable GST rates, combo deals, and discount workflows.",
            ],
            None,
        ),
        (
            "System Architecture",
            [
                "Three-tier architecture: React front-end, Node.js API layer, PostgreSQL datastore.",
                "RESTful endpoints secured via JWT, rate-limits, and role-based access control.",
                "Real-time order updates delivered using WebSockets for kitchen display.",
                "Deployment via containerized services with auto-scaling on demand.",
            ],
            None,
        ),
        (
            "Use Case Overview",
            [
                "Cashier: creates orders, accepts payments, prints/whatsapps receipts.",
                "Chef: views live order queue, updates preparation status.",
                "Manager: monitors sales KPIs, adjusts inventory, schedules staff.",
                "Customer: receives digital receipt and loyalty points snapshot.",
            ],
            None,
        ),
        (
            "Point-of-Sale Module",
            [
                "Touch-friendly menu segmented by beverages, snacks, combos, custom add-ons.",
                "Quick cart modifications, split bills, and multi-payment mode support (UPI/cash/cards).",
                "Automatic GST breakdown and ledger-friendly receipt exports.",
                "Offline caching using service workers for uninterrupted billing.",
            ],
            None,
        ),
        (
            "Inventory & Procurement Module",
            [
                "Maps recipes to raw material consumption for precise stock deductions.",
                "Batch-wise inventory with expiry tracking to reduce wastage.",
                "Supplier management for purchase orders and delivery timelines.",
                "Predictive reorder suggestions based on historical sales velocity.",
            ],
            None,
        ),
        (
            "Customer Engagement Module",
            [
                "Digital wallet and loyalty point accrual linked to phone numbers.",
                "Targeted offers via SMS/WhatsApp templates with click-through tracking.",
                "Feedback capture post-order with sentiment summarization.",
                "Heatmaps of repeat visits to drive personalized campaigns.",
            ],
            None,
        ),
        (
            "Technology Stack",
            [
                "Front-end: Next.js + TypeScript with Tailwind CSS for rapid UI composition.",
                "Back-end: Node.js (NestJS) microservices orchestrated via Express gateway.",
                "Database: Supabase/PostgreSQL with row-level security & backups.",
                "Integrations: Razorpay/UPI, Firebase Cloud Messaging, webhooks for accounting.",
            ],
            None,
        ),
        (
            "Database Design Highlights",
            [
                "Tables for menu_items, orders, order_items, payments, stocks, suppliers, customers.",
                "Use of database triggers to maintain inventory ledger and audit trails.",
                "JSONB columns to store dynamic modifiers and localized descriptions.",
                "Materialized views power dashboard KPIs and peak hour analytics.",
            ],
            None,
        ),
        (
            "Key User Interfaces",
            [
                "Dashboard shows today's revenue, average ticket size, best-selling items.",
                "Kitchen Display Screen (KDS) grouping orders by status and time elapsed.",
                "Inventory board with traffic-light indicators for stock health.",
                "Customer profile view aggregating feedback, spend, and favorites.",
            ],
            None,
        ),
        (
            "Order Workflow",
            [
                "Order initiated at POS, items selected, modifiers applied.",
                "Payment processed with automatic receipt generation and loyalty update.",
                "Order pushed to KDS; chef marks stages (accepted, preparing, ready).",
                "Completion triggers inventory deduction and analytics logging.",
            ],
            None,
        ),
        (
            "Implementation Plan",
            [
                "Phase 1 (Weeks 1-3): Requirements finalization, UI wireframes, database schema.",
                "Phase 2 (Weeks 4-7): Core POS build, order APIs, authentication modules.",
                "Phase 3 (Weeks 8-10): Inventory automation, reporting dashboards, KDS.",
                "Phase 4 (Weeks 11-12): Testing cycles, deployment, stakeholder training.",
            ],
            None,
        ),
        (
            "Core Algorithms & Logic",
            [
                "Bill computation engine handles tax slabs, discounts, rounding, and tender types.",
                "Inventory scheduler reconciles real-time sales with batch-level stock.",
                "Recommendation engine suggests combos via association rule mining.",
                "Alerting subsystem generates notifications for KPIs exceeding thresholds.",
            ],
            None,
        ),
        (
            "Security & Compliance",
            [
                "Role-based access control separating cashier, chef, manager permissions.",
                "End-to-end TLS, hashed credentials with bcrypt, and optional OTP login.",
                "Audit logs capture critical actions: voided bills, refunds, price changes.",
                "Compliant with GST invoicing rules and data residency guidelines in India.",
            ],
            None,
        ),
        (
            "Testing Strategy",
            [
                "Unit tests for billing calculations, inventory adjustments, and API endpoints.",
                "Integration tests simulating POS-KDS interactions under peak load.",
                "User acceptance testing with café staff to validate usability and workflows.",
                "Performance benchmarking ensuring sub-2 second response for 95th percentile.",
            ],
            None,
        ),
        (
            "Results & Insights",
            [
                "Pilot café reduced billing time per order by 35% after adoption.",
                "Inventory variance dropped from 11% to 3% through automated deductions.",
                "Daily dashboard enabled faster menu adjustments and improved margins.",
                "Positive feedback on bilingual UI and WhatsApp receipt sharing.",
            ],
            None,
        ),
        (
            "Cost & Feasibility Analysis",
            [
                "Development cost estimated at ₹1.8L with in-house team and open-source stack.",
                "Operational costs under ₹4k/month covering hosting, SMS, and maintenance.",
                "Break-even within 9 months for cafés averaging 150 orders/day.",
                "Scalable pricing tiers for single-outlet and franchise models.",
            ],
            None,
        ),
        (
            "Challenges & Mitigation",
            [
                "Unstable connectivity: implemented offline-first caching and sync queues.",
                "Staff adoption: provided vernacular tutorials and role-based onboarding.",
                "Data accuracy: enforced validation rules and reconciliation dashboards.",
                "Feature creep: maintained backlog with MoSCoW prioritization.",
            ],
            None,
        ),
        (
            "Future Enhancements",
            [
                "AI-driven demand forecasting and automated procurement suggestions.",
                "IoT integration for smart coffee machines and energy monitoring.",
                "Dynamic pricing experiments based on footfall and weather data.",
                "Marketplace tie-ins with delivery aggregators for omnichannel ordering.",
            ],
            None,
        ),
        (
            "Conclusion",
            [
                "Smart Cafe Management System modernizes café operations end-to-end.",
                "Provides affordable digital transformation tailored to Indian SMBs.",
                "Delivers actionable insights for profitability and customer loyalty growth.",
                "Ready roadmap for scaling into a multi-outlet SaaS platform.",
            ],
            None,
        ),
        (
            "References & Acknowledgements",
            [
                "Industry reports from NRAI India Food Services and FICCI hospitality outlook.",
                "Product benchmarks from Toast POS, Petpooja, and Loyverse documentation.",
                "Mentor guidance from Prof. Kavita Desai and Café Ananda pilot partner.",
                "Open-source communities for Next.js, Supabase, and analytics tooling.",
            ],
            None,
        ),
    ]

    for title, bullets, image in slides_data:
        add_bullet_slide(prs, title, bullets, image)

    output_path = Path("Smart Cafe Management System.pptx")
    prs.save(output_path)
    print(f"Saved presentation to {output_path.resolve()}")


if __name__ == "__main__":
    build_presentation()
